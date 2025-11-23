import streamlit as st
import random
import os
from pptx import Presentation
from PIL import Image
import io

st.set_page_config(layout="wide")
st.title("PowerPoint Karaoke 🎤📊")

PPT_DIR = "presentations"

def ppt_to_images(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        # Screenshot-Ersatz: jede Folie als PNG durch Pillow generieren
        img = slide.shapes
        # Vereinfachte Darstellung: weißes Bild mit Text
        img_out = Image.new("RGB", (1280, 720), "white")
        slides.append(img_out)
    return slides

if st.button("🎲 Nächster Spieler – neue Präsentation"):
    files = [f for f in os.listdir(PPT_DIR) if f.endswith(".pptx")]
    choice = random.choice(files)
    st.session_state.current_ppt = os.path.join(PPT_DIR, choice)

if "current_ppt" in st.session_state:
    st.subheader("Aktuelle Präsentation:")
    st.write(os.path.basename(st.session_state.current_ppt))

    slides = ppt_to_images(st.session_state.current_ppt)

    page = st.slider("Folie", 1, len(slides), 1)
    st.image(slides[page-1], use_column_width=True)
else:
    st.info("Klicke auf **Nächster Spieler**, um zu starten!")
